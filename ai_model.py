import os
import re
import json
import argparse
from datetime import datetime
from typing import List, Dict, Any
import pandas as pd
from tqdm import tqdm
from transformers import pipeline

def load_files_from_paths(paths: List[str]) -> List[str]:
    return list(paths)

def preprocess_text(text: str) -> str:
    """Preprocess text by cleaning HTML tags and normalizing whitespace"""
    # STEP 1: Remove HTML tags
    text = re.sub(r'<[^>]+>', ' ', text)
    
    # STEP 2: Replace line breaks and tabs with spaces
    text = text.replace('\n', ' ').replace('\r', ' ').replace('\t', ' ')
    
    # STEP 3: Normalize multiple whitespaces to single space
    text = re.sub(r'\s+', ' ', text)
    
    # STEP 4: Remove leading and trailing whitespace
    return text.strip()

def run_ner_pipeline(text: str, ner_pipe) -> List[Dict[str, Any]]:
    entities = ner_pipe(text)
    return entities

def is_anom_data_swl_file(file_path: str, content: str) -> bool:
    """Check if the file is an anomaly data SWL text file based on filename or content patterns"""
    filename = os.path.basename(file_path).lower()
    
    # Check filename patterns
    anom_keywords = ['anom', 'anomaly', 'swl', 'anomalous']
    if any(keyword in filename for keyword in anom_keywords):
        return True
    
    # Check content patterns (first few lines)
    content_lower = content[:1000].lower()  # Check first 1000 chars
    content_patterns = ['anomaly', 'swl', 'anomalous', 'deviation', 'outlier']
    if any(pattern in content_lower for pattern in content_patterns):
        return True
    
    return False

def process_anom_data_swl(text: str) -> Dict[str, Any]:
    """Process anomaly data SWL text files with specialized handling"""
    # STEP 1: Preprocess the raw text
    processed_text = preprocess_text(text)
    
    # STEP 2: Extract anomaly-specific patterns using regex
    anomaly_patterns = {
        'timestamps': re.findall(r'\d{4}-\d{2}-\d{2}[T\s]\d{2}:\d{2}:\d{2}', text),
        'severity_levels': re.findall(r'\b(critical|high|medium|low|warning|error|info)\b', text, re.IGNORECASE),
        'numeric_values': re.findall(r'\b\d+\.\d+\b', text),
        'swl_codes': re.findall(r'\bSWL[_-]?\w+\b', text, re.IGNORECASE),
        'anomaly_indicators': re.findall(r'\b(anomaly|anomalous|deviation|outlier|abnormal)\b', text, re.IGNORECASE)
    }
    
    # STEP 3: Build structured output with extracted information
    return {
        "raw_text": text,
        "processed_text": processed_text,
        "file_type": "anom_data_swl",
        "anomaly_patterns": anomaly_patterns,
        "extracted_info": {
            "total_anomalies": len(anomaly_patterns['anomaly_indicators']),
            "severity_distribution": {level: anomaly_patterns['severity_levels'].count(level) for level in set(anomaly_patterns['severity_levels'])},
            "timestamp_count": len(anomaly_patterns['timestamps']),
            "swl_references": len(anomaly_patterns['swl_codes'])
        }
    }

def extract_credentials(text: str) -> Dict[str, Any]:
    
    extracted_data = {
        "raw_text": text,
        "processed_text": preprocess_text(text),
        "extracted_info": {}
    }
    
    return extracted_data

def build_json_structure(file_name: str, content: Any) -> dict:
    return {
        "file_name": file_name,
        "content": content
    }

def make_json_serializable(obj):
    import numpy as np
    if isinstance(obj, dict):
        return {k: make_json_serializable(v) for k, v in obj.items()}
    elif isinstance(obj, list):
        return [make_json_serializable(v) for v in obj]
    elif isinstance(obj, (np.integer, np.int32, np.int64)):
        return int(obj)
    elif isinstance(obj, (np.floating, np.float32, np.float64)):
        return float(obj)
    elif isinstance(obj, (np.bool_)):
        return bool(obj)
    elif isinstance(obj, bytes):
        return obj.decode('utf-8', errors='ignore')
    else:
        return obj

def export_output_per_file(results: List[Dict[str, Any]], output_format: str = 'json'):
    """Export processing results to individual files in specified format"""
    import pandas as pd
    
    # STEP 1: Process each result file
    for result in results:
        input_name = result["file_name"]
        base, _ = os.path.splitext(input_name)
        content = result["content"]
        
        # STEP 2a: Export as JSON format
        if output_format == 'json':
            output_file = f"{base}.json"
            serializable_data = make_json_serializable(result)
            with open(output_file, "w", encoding="utf-8") as f:
                json.dump(serializable_data, f, ensure_ascii=False, indent=2)
            print(f"[+] Output written to {output_file}")
        
        # STEP 2b: Export as CSV format
        elif output_format == 'csv':
            output_file = f"{base}.csv"
            try:
                # STEP 2b-i: Handle structured data (lists/dicts)
                from pandas import json_normalize
                if isinstance(content, list) and all(isinstance(row, dict) for row in content):
                    df = json_normalize(content)
                    df.to_csv(output_file, index=False, encoding="utf-8")
                elif isinstance(content, dict):
                    df = json_normalize([content])
                    df.to_csv(output_file, index=False, encoding="utf-8")
                else:
                    # STEP 2b-ii: Handle plain text content
                    with open(output_file, "w", encoding="utf-8") as f:
                        f.write("file_name,content\n")
                        safe_content = str(content).replace('"', '""').replace('\n', ' ').replace('\r', '')
                        safe_name = input_name.replace('"', '""')
                        f.write(f'"{safe_name}","{safe_content}"\n')
            except Exception as e:
                # STEP 2b-iii: Handle CSV export errors
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write(f"file_name,content\n{input_name},[Error writing CSV: {e}]\n")
            print(f"[+] Output written to {output_file}")


def main(files: List[str], output_format: str = "json"):
    # STEP 1: Load and validate input files
    files = load_files_from_paths(files)
    if not files:
        print(f"[!] No supported files selected.")
        return
    
    # STEP 2: Initialize NER pipeline for entity extraction
    print("[+] Loading HuggingFace NER pipeline...")
    ner_pipe = pipeline("ner", grouped_entities=True)
    
    # STEP 3: Initialize results container
    results = []
    
    # STEP 4: Process each file based on its format
    for fpath in tqdm(files, desc="Processing files"):
        try:
            # STEP 4a: Determine file extension for format-specific processing
            ext = os.path.splitext(fpath)[1].lower()
            
            # STEP 4b: Handle CSV files
            if ext == '.csv':
                try:
                    df = pd.read_csv(fpath, dtype=str, keep_default_na=False, on_bad_lines='skip')
                    content = df.to_dict(orient='records')
                except Exception as e:
                    content = f"[Error reading CSV: {e}]"
            
            # STEP 4c: Handle Excel files
            elif ext == '.xlsx':
                try:
                    df = pd.read_excel(fpath, dtype=str)
                    content = df.fillna('').to_dict(orient='records')
                except Exception as e:
                    content = f"[Error reading XLSX: {e}]"
            
            # STEP 4d: Handle JSON files
            elif ext == '.json':
                import json as js
                try:
                    with open(fpath, 'r', encoding='utf-8', errors='ignore') as f:
                        content = js.load(f)
                except Exception as e:
                    content = f"[Error reading JSON: {e}]"
            
            # STEP 4e: Handle PowerPoint files
            elif ext == '.pptx':
                try:
                    from pptx import Presentation
                    prs = Presentation(fpath)
                    slides_text = []
                    for slide in prs.slides:
                        slide_text = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                slide_text.append(shape.text)
                        slides_text.append("\n".join(slide_text))
                    content = "\n---\n".join(slides_text)
                except Exception as e:
                    content = f"[Error reading pptx: {e}]"
            
            # STEP 4f: Handle PDF files
            elif ext == '.pdf':
                try:
                    import PyPDF2
                    with open(fpath, 'rb') as f:
                        reader = PyPDF2.PdfReader(f)
                        content = "\n".join(page.extract_text() or '' for page in reader.pages)
                except ImportError:
                    content = "[PyPDF2 not installed. Cannot extract PDF text.]"
                except Exception as e:
                    content = f"[Error reading PDF: {e}]"
            
            # STEP 4g: Handle Word documents
            elif ext in ['.doc', '.docx']:
                try:
                    import docx
                    doc = docx.Document(fpath)
                    content = "\n".join([para.text for para in doc.paragraphs])
                except ImportError:
                    content = "[python-docx not installed. Cannot extract DOCX text.]"
                except Exception as e:
                    content = f"[Error reading DOC/DOCX: {e}]"
            # STEP 4h: Handle HTML files
            elif ext in ['.html', '.htm']:
                try:
                    with open(fpath, "r", encoding="utf-8", errors="ignore") as f:
                        content = f.read()
                except Exception as e:
                    content = f"[Error reading HTML: {e}]"
            
            # STEP 4i: Handle SQL files
            elif ext == '.sql':
                try:
                    with open(fpath, "r", encoding="utf-8", errors="ignore") as f:
                        content = f.read()
                    print(f"[+] Processed SQL file: {os.path.basename(fpath)}")
                except Exception as e:
                    content = f"[Error reading SQL: {e}]"
            
            # STEP 4j: Handle DATA files
            elif ext == '.data':
                try:
                    with open(fpath, "r", encoding="utf-8", errors="ignore") as f:
                        content = f.read()
                    print(f"[+] Processed DATA file: {os.path.basename(fpath)}")
                except Exception as e:
                    content = f"[Error reading DATA: {e}]"
            
            # STEP 4k: Handle ANOM files with specialized processing
            elif ext == '.anom':
                try:
                    with open(fpath, "r", encoding="utf-8", errors="ignore") as f:
                        content = f.read()
                    # Process .anom files with specialized anomaly processing
                    processed_data = process_anom_data_swl(content)
                    content = processed_data
                    print(f"[+] Processed ANOM file with anomaly analysis: {os.path.basename(fpath)}")
                except Exception as e:
                    content = f"[Error reading ANOM: {e}]"
            
            # STEP 4l: Handle other file types with anomaly detection
            else:
                try:
                    with open(fpath, "r", encoding="utf-8", errors="ignore") as f:
                        content = f.read()
                    
                    # Check if this is an anomaly data SWL file and process accordingly
                    if is_anom_data_swl_file(fpath, content):
                        processed_data = process_anom_data_swl(content)
                        content = processed_data
                        print(f"[+] Processed as anomaly data SWL file: {os.path.basename(fpath)}")
                        
                except Exception as e:
                    # STEP 4m: Handle binary files as fallback
                    try:
                        with open(fpath, "rb") as f:
                            content = f"[Binary file, {len(f.read())} bytes]"
                    except Exception as e2:
                        content = f"[Could not read file: {e2}]"
            
            # STEP 5: Build JSON structure for each processed file
            result = build_json_structure(
                file_name=os.path.basename(fpath),
                content=content
            )
            results.append(result)
        except Exception as e:
            print(f"[!] Error processing {fpath}: {e}")
    
    # STEP 6: Export results to specified output format
    export_output_per_file(results, output_format)


if __name__ == "__main__":
    import sys
    import tkinter as tk
    from tkinter import filedialog

    root = tk.Tk()
    root.withdraw()
    root.update()
    file_paths = filedialog.askopenfilenames(title="Select files to process", filetypes=[
        ("All Supported", ".txt .log .csv .xlsx .json .pptx .pdf .doc .docx .html .htm .sql .data .anom"),
        ("Text Files", ".txt"),
        ("SQL Files", ".sql"),
        ("Data Files", ".data"),
        ("Anomaly Files", ".anom"),
        ("Log Files", ".log"),
        ("CSV Files", ".csv"),
        ("Excel Files", ".xlsx"),
        ("JSON Files", ".json"),
        ("PowerPoint Files", ".pptx"),
        ("PDF Files", ".pdf"),
        ("Word Files", ".doc .docx"),
        ("HTML Files", ".html .htm"),
        ("All Files", "*.*")
    ])
    root.destroy()
    file_paths = list(file_paths)
    if not file_paths:
        print("[!] No files selected. Exiting.")
        sys.exit(1)


    while True:
        out_fmt = input("Choose output format (json/csv): ").strip().lower()
        if out_fmt in ("json", "csv"):
            break
        print("Invalid choice. Please enter 'json' or 'csv'.")

    main(file_paths, out_fmt)

