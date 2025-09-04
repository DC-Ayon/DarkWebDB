import os
import re
import json
import argparse
import zipfile
import tarfile
import tempfile
import shutil
import subprocess
import sys
from datetime import datetime
from typing import List, Dict, Any

def safe_import(package_name, import_name=None):
    import_name = import_name or package_name
    try:
        return __import__(import_name)
    except ImportError:
        print(f"[!] Installing missing package: {package_name}")
        subprocess.check_call([sys.executable, "-m", "pip", "install", package_name])
        return __import__(import_name)

# Import external modules with auto-install
pd = safe_import("pandas")
patoolib = safe_import("patool", "patoolib")
from tqdm import tqdm
from transformers import pipeline

def is_supported_file(file_path: str) -> bool:
    supported_extensions = [
        ".csv", ".xlsx", ".xls", ".json", ".txt", ".pptx", ".pdf", ".doc", ".docx", ".html", ".htm", ".sql", ".data", ".anom"
    ]
    return any(file_path.lower().endswith(ext) for ext in supported_extensions)

def is_archive(file_path: str) -> bool:
    archive_exts = [".zip", ".tar", ".tar.gz", ".tgz", ".rar", ".7z"]
    return any(file_path.lower().endswith(ext) for ext in archive_exts)

def extract_archive(file_path: str) -> str:
    temp_dir = tempfile.mkdtemp()
    try:
        if file_path.lower().endswith(".zip"):
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                zip_ref.extractall(temp_dir)
        elif file_path.lower().endswith((".tar", ".tar.gz", ".tgz")):
            with tarfile.open(file_path, 'r:*') as tar_ref:
                tar_ref.extractall(temp_dir)
        else:
            patoolib.extract_archive(file_path, outdir=temp_dir, verbosity=-1)
    except Exception as e:
        print(f"[✗] Failed to extract archive {file_path}: {e}")
        shutil.rmtree(temp_dir)
        return None
    return temp_dir

def try_encodings(file_path, encodings=["utf-8", "latin1", "windows-1252", "ISO-8859-1"]) -> str:
    for enc in encodings:
        try:
            with open(file_path, "r", encoding=enc, errors='strict') as f:
                return f.read()
        except UnicodeDecodeError:
            continue
    print(f"[✗] Could not decode {file_path} with known encodings.")
    return None

def load_files_from_paths(paths: List[str]) -> List[str]:
    return list(paths)

def preprocess_text(text: str) -> str:
    """Preprocess text by cleaning HTML tags and normalizing whitespace"""
    # STEP 1: Remove HTML tags
    text = re.sub(r'<[^>]+>', ' ', text)
    
    # STEP 2: Replace line breaks and tabs with spaces
    text = text.replace('\n', ' ').replace('\r', ' ').replace('\t', ' ')
    
    # STEP 3: Normalize multiple whitespaces to single space
    text = re.sub(r'\s+', ' ', text)
    
    # STEP 4: Remove leading and trailing whitespace
    return text.strip()

def run_ner_pipeline(text: str, ner_pipe) -> List[Dict[str, Any]]:
    entities = ner_pipe(text)
    return entities

def is_anom_data_swl_file(file_path: str, content: str) -> bool:
    """Check if the file is an anomaly data SWL text file based on filename or content patterns"""
    filename = os.path.basename(file_path).lower()
    
    # Check filename patterns
    anom_keywords = ['anom', 'anomaly', 'swl', 'anomalous']
    if any(keyword in filename for keyword in anom_keywords):
        return True
    
    # Check content patterns (first few lines)
    content_lower = content[:1000].lower()  # Check first 1000 chars
    content_patterns = ['anomaly', 'swl', 'anomalous', 'deviation', 'outlier']
    if any(pattern in content_lower for pattern in content_patterns):
        return True
    
    return False

def process_anom_data_swl(text: str) -> Dict[str, Any]:
    """Process anomaly data SWL text files with specialized handling"""
    # STEP 1: Preprocess the raw text
    processed_text = preprocess_text(text)
    
    # STEP 2: Extract anomaly-specific patterns using regex
    anomaly_patterns = {
        'timestamps': re.findall(r'\d{4}-\d{2}-\d{2}[T\s]\d{2}:\d{2}:\d{2}', text),
        'severity_levels': re.findall(r'\b(critical|high|medium|low|warning|error|info)\b', text, re.IGNORECASE),
        'numeric_values': re.findall(r'\b\d+\.\d+\b', text),
        'swl_codes': re.findall(r'\bSWL[_-]?\w+\b', text, re.IGNORECASE),
        'anomaly_indicators': re.findall(r'\b(anomaly|anomalous|deviation|outlier|abnormal)\b', text, re.IGNORECASE)
    }
    
    # STEP 3: Build structured output with extracted information
    return {
        "raw_text": text,
        "processed_text": processed_text,
        "file_type": "anom_data_swl",
        "anomaly_patterns": anomaly_patterns,
        "extracted_info": {
            "total_anomalies": len(anomaly_patterns['anomaly_indicators']),
            "severity_distribution": {level: anomaly_patterns['severity_levels'].count(level) for level in set(anomaly_patterns['severity_levels'])},
            "timestamp_count": len(anomaly_patterns['timestamps']),
            "swl_references": len(anomaly_patterns['swl_codes'])
        }
    }

def extract_credentials(text: str) -> Dict[str, Any]:
    
    extracted_data = {
        "raw_text": text,
        "processed_text": preprocess_text(text),
        "extracted_info": {}
    }
    
    return extracted_data

def build_json_structure(file_name: str, content: Any) -> dict:
    return {
        "file_name": file_name,
        "content": content
    }

def make_json_serializable(obj):
    import numpy as np
    if isinstance(obj, dict):
        return {k: make_json_serializable(v) for k, v in obj.items()}
    elif isinstance(obj, list):
        return [make_json_serializable(v) for v in obj]
    elif isinstance(obj, (np.integer, np.int32, np.int64)):
        return int(obj)
    elif isinstance(obj, (np.floating, np.float32, np.float64)):
        return float(obj)
    elif isinstance(obj, (np.bool_)):
        return bool(obj)
    elif isinstance(obj, bytes):
        return obj.decode('utf-8', errors='ignore')
    else:
        return obj

def export_output_per_file(results: List[Dict[str, Any]], output_format: str = 'json', output_dir: str = None):
    """Export processing results to individual files in specified format and output directory"""
    import pandas as pd
    import os
    output_dir = output_dir or os.getcwd()
    
    # STEP 1: Process each result file
    for result in results:
        input_name = result["file_name"]
        base, _ = os.path.splitext(input_name)
        content = result["content"]
        
        # Compose output path
        if output_format == 'json':
            output_file = os.path.join(output_dir, f"{base}.json")
            serializable_data = make_json_serializable(result)
            with open(output_file, "w", encoding="utf-8") as f:
                json.dump(serializable_data, f, ensure_ascii=False, indent=2)
            print(f"[+] Output written to {output_file}")
        elif output_format == 'csv':
            output_file = os.path.join(output_dir, f"{base}.csv")
            try:
                from pandas import json_normalize
                if isinstance(content, list) and all(isinstance(row, dict) for row in content):
                    df = json_normalize(content)
                    df.to_csv(output_file, index=False, encoding="utf-8")
                elif isinstance(content, dict):
                    df = json_normalize([content])
                    df.to_csv(output_file, index=False, encoding="utf-8")
                else:
                    with open(output_file, "w", encoding="utf-8") as f:
                        f.write("file_name,content\n")
                        safe_content = str(content).replace('"', '""').replace('\n', ' ').replace('\r', '')
                        safe_name = input_name.replace('"', '""')
                        f.write(f'"{safe_name}","{safe_content}"\n')
            except Exception as e:
                # STEP 2b-iii: Handle CSV export errors
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write(f"file_name,content\n{input_name},[Error writing CSV: {e}]\n")
            print(f"[+] Output written to {output_file}")


def main(files: List[str], output_format: str = "json", output_dir: str = None):
    # STEP 1: Load and validate input files
    files = load_files_from_paths(files)
    if not files:
        print(f"[!] No supported files selected.")
        return

    # STEP 2: Expand archives and collect all files
    all_files = []
    for f in files:
        if is_archive(f):
            extracted_dir = extract_archive(f)
            if extracted_dir:
                for root, _, filenames in os.walk(extracted_dir):
                    for name in filenames:
                        full_path = os.path.join(root, name)
                        if is_supported_file(full_path):
                            all_files.append(full_path)
        elif is_supported_file(f):
            all_files.append(f)

    if not all_files:
        print("[!] No supported files found after archive extraction.")
        return

    # STEP 3: Initialize NER pipeline for entity extraction
    print("[+] Loading HuggingFace NER pipeline...")
    ner_pipe = pipeline("ner", grouped_entities=True)

    # STEP 4: Initialize results container
    results = []

    # STEP 5: Process each file based on its format
    for fpath in tqdm(all_files, desc="Processing files"):
        try:
            ext = os.path.splitext(fpath)[1].lower()
            content = None
            # CSV
            if ext == '.csv':
                try:
                    df = pd.read_csv(fpath, dtype=str, keep_default_na=False, on_bad_lines='skip')
                    content = df.to_dict(orient='records')
                except Exception as e:
                    content = f"[Error reading CSV: {e}]"
            # Excel
            elif ext in ('.xlsx', '.xls'):
                try:
                    df = pd.read_excel(fpath, dtype=str)
                    content = df.fillna('').to_dict(orient='records')
                except Exception as e:
                    content = f"[Error reading Excel: {e}]"
            # JSON
            elif ext == '.json':
                try:
                    with open(fpath, 'r', encoding='utf-8', errors='ignore') as f:
                        content = json.load(f)
                except Exception as e:
                    content = f"[Error reading JSON: {e}]"
            # TXT
            elif ext == '.txt':
                txt_content = try_encodings(fpath)
                content = txt_content if txt_content else ""
            # PPTX
            elif ext == '.pptx':
                try:
                    from pptx import Presentation
                    prs = Presentation(fpath)
                    slides_text = []
                    for slide in prs.slides:
                        slide_text = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                slide_text.append(shape.text)
                        slides_text.append("\n".join(slide_text))
                    content = "\n---\n".join(slides_text)
                except Exception as e:
                    content = f"[Error reading pptx: {e}]"
            # PDF
            elif ext == '.pdf':
                try:
                    import PyPDF2
                    with open(fpath, 'rb') as f:
                        reader = PyPDF2.PdfReader(f)
                        content = "\n".join(page.extract_text() or '' for page in reader.pages)
                except ImportError:
                    content = "[PyPDF2 not installed. Cannot extract PDF text.]"
                except Exception as e:
                    content = f"[Error reading PDF: {e}]"
            # DOC/DOCX
            elif ext in ['.doc', '.docx']:
                try:
                    import docx
                    doc = docx.Document(fpath)
                    content = "\n".join([para.text for para in doc.paragraphs])
                except ImportError:
                    content = "[python-docx not installed. Cannot extract DOCX text.]"
                except Exception as e:
                    content = f"[Error reading DOC/DOCX: {e}]"
            # HTML
            elif ext in ['.html', '.htm']:
                try:
                    with open(fpath, "r", encoding="utf-8", errors="ignore") as f:
                        content = f.read()
                except Exception as e:
                    content = f"[Error reading HTML: {e}]"
            # SQL
            elif ext == '.sql':
                try:
                    with open(fpath, "r", encoding="utf-8", errors="ignore") as f:
                        content = f.read()
                    print(f"[+] Processed SQL file: {os.path.basename(fpath)}")
                except Exception as e:
                    content = f"[Error reading SQL: {e}]"
            # DATA
            elif ext == '.data':
                try:
                    with open(fpath, "r", encoding="utf-8", errors="ignore") as f:
                        content = f.read()
                    print(f"[+] Processed DATA file: {os.path.basename(fpath)}")
                except Exception as e:
                    content = f"[Error reading DATA: {e}]"
            # ANOM
            elif ext == '.anom':
                try:
                    with open(fpath, "r", encoding="utf-8", errors="ignore") as f:
                        content = f.read()
                    processed_data = process_anom_data_swl(content)
                    content = processed_data
                    print(f"[+] Processed ANOM file with anomaly analysis: {os.path.basename(fpath)}")
                except Exception as e:
                    content = f"[Error reading ANOM: {e}]"
            # Fallback/Other
            else:
                try:
                    with open(fpath, "r", encoding="utf-8", errors="ignore") as f:
                        content = f.read()
                    if is_anom_data_swl_file(fpath, content):
                        processed_data = process_anom_data_swl(content)
                        content = processed_data
                        print(f"[+] Processed as anomaly data SWL file: {os.path.basename(fpath)}")
                except Exception as e:
                    try:
                        with open(fpath, "rb") as f:
                            content = f"[Binary file, {len(f.read())} bytes]"
                    except Exception as e2:
                        content = f"[Could not read file: {e2}]"
            # NLP/NER pipeline integration
            if isinstance(content, str) and content:
                entities = run_ner_pipeline(content, ner_pipe)
                result = build_json_structure(fpath, {"entities": entities, "raw_text": content})
            else:
                result = build_json_structure(fpath, content)
            results.append(result)
        except Exception as e:
            print(f"[!] Error processing {fpath}: {e}")
    export_output_per_file(results, output_format, output_dir)


if __name__ == "__main__":
    import sys
    import tkinter as tk
    from tkinter import filedialog

    # CLI/GUI logic
    root = tk.Tk()
    root.withdraw()
    print("Select files to process...")
    SUPPORTED_EXTS = [
        ("All Supported Files", "*.csv *.xlsx *.xls *.json *.txt *.pptx *.pdf *.doc *.docx *.html *.htm *.sql *.data *.anom *.zip *.tar *.tar.gz *.tgz *.rar *.7z"),
        ("CSV Files", "*.csv"),
        ("Excel Files", "*.xlsx;*.xls"),
        ("JSON Files", "*.json"),
        ("Text Files", "*.txt"),
        ("PowerPoint Files", "*.pptx"),
        ("PDF Files", "*.pdf"),
        ("Word Files", "*.doc;*.docx"),
        ("HTML Files", "*.html;*.htm"),
        ("SQL Files", "*.sql"),
        ("Data Files", "*.data"),
        ("Anomaly Files", "*.anom"),
        ("Archive Files", "*.zip;*.tar;*.tar.gz;*.tgz;*.rar;*.7z"),
        ("HTML Files", ".html .htm"),
        ("All Files", "*.*")
    ]
    file_paths = filedialog.askopenfilenames(
        title="Select files for processing",
        filetypes=SUPPORTED_EXTS
    )
    root.destroy()
    file_paths = list(file_paths)
    if not file_paths:
        print("[!] No files selected. Exiting.")
        sys.exit(0)

    output_format = input("Choose output format (json/csv): ").strip().lower()
    while output_format not in ["json", "csv"]:
        print("Invalid choice. Please enter 'json' or 'csv'.")
        output_format = input("Choose output format (json/csv): ").strip().lower()

    print("Select output folder for results...")
    output_dir = filedialog.askdirectory(title="Select output folder for results")
    if not output_dir:
        print("[!] No output folder selected. Exiting.")
        sys.exit(0)

    main(file_paths, output_format, output_dir)
